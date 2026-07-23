import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    NAVY = RGBColor(15, 34, 64)
    BLUE = RGBColor(0, 114, 206)
    LIGHT_BG = RGBColor(245, 247, 250)
    DARK_TEXT = RGBColor(30, 30, 30)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 100, 100)
    ACCENT_ORANGE = RGBColor(227, 85, 33)
    CARD_BG = RGBColor(235, 240, 248)
    CARD_BORDER = RGBColor(200, 215, 235)
    WARN_BG = RGBColor(254, 243, 235)
    WARN_BORDER = RGBColor(240, 150, 100)

    def set_slide_background(slide, color=LIGHT_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.color.rgb = color
        return bg

    def add_header(slide, title_text, category_text="AZURE CYCLECLOUD MSP TRAINING & KT STATUS"):
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = NAVY
        top_bar.line.color.rgb = NAVY
        
        tf = top_bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.6)
        tf.margin_top = Inches(0.15)
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = BLUE
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title, items, bg_color=CARD_BG, border_color=CARD_BORDER, title_color=NAVY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}" if not item.startswith("  ") else item
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, NAVY)
    
    tf1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0)).text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "AZURE CYCLECLOUD"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf1.add_paragraph()
    p.text = "MSP 운영 교육 & KT 현황 브리핑"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p = tf1.add_paragraph()
    p.text = "최초 사이트/클러스터 구축, 오케스트레이션, 실습 및 KT 환경 지침"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 215, 235)
    
    p = tf1.add_paragraph()
    p.text = "\n발표/지원: Azure SRE Team (이동훈, 김태진, 지민근) | 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 2: Objectives & Agenda
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "교육 목표 및 커리큘럼 (Agenda)")

    add_card(s2, 0.6, 1.4, 5.8, 5.5, "🎯 교육 목표 (Training Objectives)", [
        "Azure CycleCloud 최초 포털 및 클러스터 구축 절차 마스터",
        "웹 포털 및 CLI(cyclecloud/azslurm)를 통한 클러스터 제어 습득",
        "KT CycleCloud 운영 특이사항 및 제약 조건 파악",
        "  - RI 적용 GPU VM scale-in 방지 및 노드 재시작 최소화",
        "  - cloud-init vs cluster-init 구성 지침",
        "  - Slurm Job Accounting (MySQL Flexible Server)",
        "  - Managed Prometheus + Grafana GPU 모니터링",
        "실습 환경(Bicep IaC) 기반 MSP 전용 운영/트러블슈팅 역량 확보"
    ])

    add_card(s2, 6.8, 1.4, 5.9, 5.5, "📋 교육 아젠다 (Agenda)", [
        "1부. Azure CycleCloud & Slurm 개요",
        "  - 개념, 아키텍처, 최초 클러스터 구축 및 Slurm CLI",
        "2부. KT CycleCloud 구축 현황 & 핵심 운영 지침",
        "  - Mexico 리전 이관 현황, 버전 체계, CLI 서버 접속",
        "  - GPU VM RI 규정, cluster-init, Job Accounting",
        "  - 사용자 관리, 스토리지 마운트, GPU 모니터링",
        "3부. MSP 핸즈온 실습 런북 (Lab Guide)",
        "  - 12개 실습 모듈 및 트러블슈팅 로그 진단 절차"
    ])

    # -------------------------------------------------------------
    # SLIDE 3: CycleCloud Architecture
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "1부. CycleCloud & Slurm 구성 아키텍처")

    add_card(s3, 0.6, 1.4, 3.8, 5.5, "① Control Plane (Server)", [
        "CycleCloud Server VM (cc-server)",
        "웹 포털 GUI (HTTPS 443)",
        "cycle_server 서비스 & CLI",
        "System-assigned Managed Identity",
        "  - 구독 Contributor 권한 부여",
        "  - 자격증명/암호 없이 노드 프로비저닝"
    ])

    add_card(s3, 4.7, 1.4, 3.8, 5.5, "② Scheduler & Compute Nodes", [
        "Slurm Master (Scheduler 노드)",
        "  - 워크로드 스케줄링 & 자원 할당",
        "  - Slurmctld 데몬 기동",
        "Execute (Worker) 노드",
        "  - VMSS 기반 필요 시 자동 생성",
        "  - Job 완료 후 유휴 시 자동 감설"
    ])

    add_card(s3, 8.8, 1.4, 3.9, 5.5, "③ Storage & Locker", [
        "/shared, /sched (NFS)",
        "  - 모든 노드 공용 공유 저장소",
        "  - Slurm 메타데이터 및 스케줄러 설정",
        "Blob Storage (Locker)",
        "  - 스토리지 계정 (cclkekwphusd3i)",
        "  - cluster-init 프로젝트, 템플릿 보관"
    ])

    # -------------------------------------------------------------
    # SLIDE 4: First-Time Cluster Setup & Initialization Guide (NEW!)
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "1부. 최초 클러스터 구축 단계별 가이드 (First-Time Setup)")

    add_card(s4, 0.6, 1.4, 3.8, 5.5, "1단계: Site 초기화 (Wizard)", [
        "1. 포털 접속: https://<server-ip>",
        "2. 약관 동의 & 관리자 계정 생성",
        "3. Azure 구독 (Credential) 등록:",
        "  - Service Type: Managed Identity",
        "  - Subscription ID 입력",
        "  - Default Region: Korea Central",
        "  - Resource Group: rg-cyclecloud-training",
        "  - Locker: cclkekwphusd3i"
    ])

    add_card(s4, 4.7, 1.4, 3.8, 5.5, "2단계: Slurm 템플릿 구성", [
        "1. Clusters > '+' (New Cluster)",
        "2. 스케줄러 템플릿: Slurm 선택",
        "3. 파라미터 설정:",
        "  - Cluster Name: slurm-first-cluster",
        "  - Scheduler VM: Standard_D4s_v5",
        "  - HPC Worker VM: Standard_D4s_v5",
        "  - Max Cores: 100 (Autoscale)",
        "  - Subnet: compute (10.0.1.0/24)"
    ])

    add_card(s4, 8.8, 1.4, 3.9, 5.5, "3단계: 최초 기동 & 검증", [
        "1. Start 버튼 클릭 (노드 상태 변화):",
        "  - Off -> Acquiring -> Preparing -> Ready",
        "2. 스케줄러 노드 SSH 접속:",
        "  - cyclecloud connect scheduler",
        "3. 상태 조회 & 테스트 작업 제출:",
        "  - sinfo / squeue",
        "  - srun -N 1 hostname",
        "  - sbatch test_job.sh"
    ])

    # -------------------------------------------------------------
    # SLIDE 5: Slurm Commands & CLI
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "1부. Slurm 주요 명령어 & CycleCloud CLI")

    add_card(s5, 0.6, 1.4, 5.8, 5.5, "⚡ Slurm 주요 명령어 요약", [
        "sinfo : 파티션 및 노드 상태 조회 (idle, alloc, down, drain)",
        "squeue : Slurm 큐의 Job 상태 및 대기 사유(Reason) 확인",
        "sbatch <script.sh> : 비동기 배치 작업 스크립트 제출",
        "srun <cmd> : Slurm 자원을 즉시 할당받아 대화형 실행",
        "sacct : 완료된 Job 실행 이력 및 CPU/GPU/메모리 사용량 조회",
        "scontrol show job <jobid> : Job 상세 정보 및 대기 원인 분석",
        "scontrol reconfigure : Slurm 설정(azure.conf) 즉시 재적용"
    ])

    add_card(s5, 6.8, 1.4, 5.9, 5.5, "🛠️ CycleCloud CLI (`cyclecloud` & `azslurm`)", [
        "CycleCloud Master Server에 사전 설치 완료",
        "cyclecloud show_cluster <cluster> : 클러스터 상태 확인",
        "cyclecloud show_nodes -c <cluster> : 노드별 상세 상태 확인",
        "cyclecloud start_cluster / terminate_cluster : 클러스터 제어",
        "cyclecloud project upload <locker> : cluster-init 업로드",
        "azslurm resume --node-list <node> : 수동 노드 프로비저닝",
        "azslurm suspend --node-list <node> : 수동 노드 회수/삭제"
    ])

    # -------------------------------------------------------------
    # SLIDE 6: KT Status Overview
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "2부. KT CycleCloud 현재 구축 현황 (2026년 기준)")

    add_card(s6, 0.6, 1.4, 3.8, 5.5, "📍 리전 이관 현황", [
        "기존 구성: UAEN, KRC, KRS 3개 리전 분산 운영",
        "2026년 6월: Mexico Region으로 일괄 통합 이관",
        "현재 규모:",
        "  - 운영 환경: Slurm 클러스터 1개 (노드 약 100개)",
        "  - 개발 환경: Single VM 각각 4개"
    ])

    add_card(s6, 4.7, 1.4, 3.8, 5.5, "🔍 버전 체계 관리", [
        "기술 지원 시 3가지 버전 정보 확인 필수:",
        "  1. CycleCloud Server 버전 (UI 우측 상단 '?')",
        "  2. Slurm Version (Cluster Edit > Advanced)",
        "  3. Jetpack / azslurm 버전 (Master 노드 CLI)",
        "버전 확인 가이드 및 스크립트 문서화 완료"
    ])

    add_card(s6, 8.8, 1.4, 3.9, 5.5, "💻 CLI & 접근 제어", [
        "CycleCloud Master Server (cc-server)에 설치",
        "CycleCloud GUI 접속도 서버 공인 IP로 처리",
        "Azure Portal '명령 실행(Run Command)' 활용:",
        "  - SSH 키 없이도 Azure RBAC 권한만으로 서버 접근 및 로그 확인 가능"
    ])

    # -------------------------------------------------------------
    # SLIDE 7: KT Guideline 1 - GPU RI & Reboot Policy
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "2부. KT 운영 지침 ①: GPU VM RI & 노드 재시작 최소화")

    add_card(s7, 0.6, 1.4, 12.133, 2.3, "⚠️ [최우선 지침] 노드 재시작 최소화 원칙", [
        "CycleCloud 지원 및 장애 조치 시 노드 재시작을 최소화하는 것이 가장 중요합니다.",
        "원인: GPU VM Capa(용량) 부족으로 인해 노드 재시작/재생성 시 재할당(Acquiring)이 실패할 수 있음",
        "조치: 재시작이 불가피한 경우 반드시 Azure Capa 팀과 사전 일정을 조율한 후 진행하십시오."
    ], bg_color=WARN_BG, border_color=WARN_BORDER, title_color=ACCENT_ORANGE)

    add_card(s7, 0.6, 3.9, 5.8, 3.0, "💳 Reserved Instance (RI) 적용 규정", [
        "KT CycleCloud GPU VM들은 Reserved Instance(RI)가 적용되어 있습니다.",
        "의도치 않은 Scale-in으로 인해 RI 자원이 놀거나 회수되는 상황을 방지해야 합니다."
    ])

    add_card(s7, 6.8, 3.9, 5.9, 3.0, "🛡️ Scale-in 방지 설정 (slurm.conf)", [
        "Slurm 기본 Suspend time은 300초(5분)로 유휴 시 자동 감설됨",
        "UI의 'keep_alive' 옵션 오작동 사례 존재",
        "해결: /etc/slurm/slurm.conf 에 SuspendExcParts=hpc 설정 적용"
    ])

    # -------------------------------------------------------------
    # SLIDE 8: KT Guideline 2 - cloud-init vs cluster-init
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "2부. KT 운영 지침 ②: cloud-init vs cluster-init")

    add_card(s8, 0.6, 1.4, 5.8, 5.5, "❌ cloud-init 사용 시 주의사항 (금지)", [
        "cloud-init은 VMSS level user-data 스크립트입니다.",
        "수정 후 노드 추가/재기동 시 아래 오류 발생하며 노드 할당 실패:",
        "  'This node does not match existing scaleset attribute: CloudInit'",
        "원인: 기존 VMSS 속성과 신규 속성 불일치",
        "위험: 문제 해소를 위해 전체 클러스터 재기동 필요"
    ], bg_color=WARN_BG, border_color=WARN_BORDER, title_color=ACCENT_ORANGE)

    add_card(s8, 6.8, 1.4, 5.9, 5.5, "✅ cluster-init 권장 및 KT 구성 현황", [
        "cluster-init은 CycleCloud 전용 수렴(converge) 프로젝트 방식",
        "개별 노드/배열 단위로 변경 사항 안전하게 적용 가능",
        "KT 구성 현황:",
        "  - Blob Storage (Locker: cclkekwphusd3i)에 4~5개 스크립트 저장",
        "  - 주요 스크립트: Scale-in 방지, NVIDIA 드라이버 버전 업, 마운트 등",
        "  - jetpack converge 로 노드 재부팅 없이 스크립트 재실행 가능"
    ])

    # -------------------------------------------------------------
    # SLIDE 9: KT Guideline 3 - Accounting & User Management
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "2부. KT 운영 지침 ③: Job Accounting & 사용자 관리")

    add_card(s9, 0.6, 1.4, 5.8, 5.5, "📊 Job Accounting (sacct) 연동", [
        "Job 실행 이력 및 CPU/GPU/메모리 사용량을 영구 보존",
        "DB 구성: Azure Database for MySQL - Flexible Server",
        "  - VNet 내 Private Subnet 배치로 보안 강화",
        "slurmdbd 데몬 설정:",
        "  - AccountingStorageType=accounting_storage/slurmdbd",
        "  - /etc/slurm/AzureCA.pem SSL 인증서 다운로드 및 연결",
        "  - sacct 및 sreport 명령으로 유저별/프로젝트별 정산 가능"
    ])

    add_card(s9, 6.8, 1.4, 5.9, 5.5, "👤 사용자 및 계정 관리 (User Management)", [
        "CycleCloud 사용자 관리 방식: Built-in 방식 사용 중",
        "사용자 생성 및 권한 부여:",
        "  - CycleCloud GUI > Users 탭에서 계정 생성",
        "  - Node Settings: Keypair (pem 파일) 등록으로 노드 SSH 접속",
        "  - Cluster에 Users 추가 시 jetpack 데몬이 노드 OS 계정 자동 생성",
        "  - Sudo 권한 필요 시 sudoers 설정 반영"
    ])

    # -------------------------------------------------------------
    # SLIDE 10: KT Guideline 4 - Storage & GPU Monitoring
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "2부. KT 운영 지침 ④: 스토리지 & GPU 모니터링")

    add_card(s10, 0.6, 1.4, 5.8, 5.5, "💾 스토리지 & 디스크 마운트 구성", [
        "OS Disk: 노드별 OS Disk 1TB 설정",
        "  - (주의: BootDiskSize 변경 시 노드 전체 Terminate/Start 필요)",
        "공유 스토리지 마운트:",
        "  - Blob Storage NFS : 2개 마운트 (데이터 저장)",
        "  - Azure Managed Lustre : 512TB 초고속 고성능 마운트",
        "Locker Storage: cluster-init 및 클러스터 메타데이터 보관"
    ])

    add_card(s10, 6.8, 1.4, 5.9, 5.5, "📈 GPU 모니터링 (Prometheus + Grafana)", [
        "아키텍처: Node Exporter -> Prometheus -> Azure Managed Prometheus (remote_write) -> Managed Grafana 시각화",
        "Azure Managed Workspace 활용하여 실시간 GPU 사용률/메모리 감시",
        "Mexico Region 특이사항:",
        "  - Mexico Region은 Managed Grafana 미지원이므로 타 리전 연결 또는 Self-hosted Grafana 설치 필요"
    ])

    # -------------------------------------------------------------
    # SLIDE 11: MSP Lab Roadmap Overview
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "3부. MSP 핸즈온 실습 모듈 목차 (Lab Guide)")

    add_card(s11, 0.6, 1.4, 3.8, 5.5, "📚 기본 과정 (Module 01~05)", [
        "01. 환경 개요 및 아키텍처",
        "02. CycleCloud 포털 사용 방법",
        "03. 최초 클러스터 구축 가이드",
        "04. 노드 증/감설 및 사이즈 변경",
        "05. 스토리지 계정 & 디스크 마운트"
    ])

    add_card(s11, 4.7, 1.4, 3.8, 5.5, "🚀 심화 과정 (Module 06~10)", [
        "06. cluster-init & 커스텀 스크립트",
        "07. Slurm Job Accounting 구축",
        "08. Built-in 사용자 및 키페어 관리",
        "09. 파티션 관리 및 추가 설정",
        "10. Azure Managed GPU 모니터링"
    ])

    add_card(s11, 8.8, 1.4, 3.9, 5.5, "🔧 운용 및 런북 (Module 11~12)", [
        "11. 트러블슈팅 & 로그 진단",
        "  - cycle_server.log 분석",
        "  - jetpack.log & cloud-init 진단",
        "  - az vm run-command 활용",
        "12. 데모 런북 (진행자 체크리스트)"
    ])

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    build_presentation("docs/Cyclecloud_MSP_Training.pptx")
    try:
        build_presentation("docs/Cyclecloud_260722_Updated.pptx")
    except Exception as e:
        print(f"Note: Could not overwrite Cyclecloud_260722.pptx: {e}")
