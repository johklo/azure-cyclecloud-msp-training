"""Cloud MSP 대상 Azure CycleCloud 운영 소개 프레젠테이션 생성 스크립트.

대상: 고객사(예: KT) CycleCloud를 대신 운영·지원하는 Cloud MSP 운영/SRE 담당자.
관점: 'MSP가 무엇을, 어떻게 운영·지원해야 하는가' (가치제안이 아닌 운영 온보딩).
근간: docs 폴더의 교육 매뉴얼(01~12) 및 KT 운영 브리핑 자료.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Color Palette (교육 덱과 동일 톤) ---
NAVY = RGBColor(15, 34, 64)
BLUE = RGBColor(0, 114, 206)
LIGHT_BG = RGBColor(245, 247, 250)
DARK_TEXT = RGBColor(30, 30, 30)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
ACCENT_ORANGE = RGBColor(227, 85, 33)
CARD_BG = RGBColor(235, 240, 248)
CARD_BORDER = RGBColor(200, 215, 235)
WARN_BG = RGBColor(254, 243, 235)
WARN_BORDER = RGBColor(240, 150, 100)
GREEN_BG = RGBColor(233, 245, 236)
GREEN_BORDER = RGBColor(150, 200, 160)
GREEN_TITLE = RGBColor(30, 120, 60)

CATEGORY = "AZURE CYCLECLOUD | CLOUD MSP 운영 온보딩 (OPERATIONS ONBOARDING)"


def build(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, color=LIGHT_BG):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.color.rgb = color
        return s

    def header(slide, title, category=CATEGORY):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.color.rgb = NAVY
        tf = bar.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.15)
        p0 = tf.paragraphs[0]; p0.text = category.upper()
        p0.font.size = Pt(10); p0.font.bold = True; p0.font.color.rgb = BLUE
        p1 = tf.add_paragraph(); p1.text = title
        p1.font.size = Pt(22); p1.font.bold = True; p1.font.color.rgb = WHITE

    def card(slide, left, top, width, height, title, items,
             bg_color=CARD_BG, border_color=CARD_BORDER, title_color=NAVY, title_size=16):
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
        c.fill.solid(); c.fill.fore_color.rgb = bg_color
        c.line.color.rgb = border_color; c.line.width = Pt(1.5)
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
        p0 = tf.paragraphs[0]; p0.text = title
        p0.font.size = Pt(title_size); p0.font.bold = True; p0.font.color.rgb = title_color
        for item in items:
            p = tf.add_paragraph()
            p.text = item if item.startswith("  ") else f"• {item}"
            p.font.size = Pt(12); p.font.color.rgb = DARK_TEXT; p.space_before = Pt(4)

    def textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame

    # ---------------------------------------------------------------
    # SLIDE 1 : Title
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    tf = textbox(s, 1.0, 2.0, 11.333, 4.2); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "AZURE CYCLECLOUD"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = BLUE
    p = tf.add_paragraph(); p.text = "Cloud MSP 운영 온보딩 가이드"
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "고객 HPC 클러스터를 대신 운영·지원하기 위해 MSP가 알아야 할 핵심"
    p.font.size = Pt(18); p.font.color.rgb = RGBColor(200, 215, 235)
    p = tf.add_paragraph(); p.text = "\n발표/지원: Azure SRE Team (이동훈, 김태진, 지민근) | 2026"
    p.font.size = Pt(14); p.font.color.rgb = GRAY

    # ---------------------------------------------------------------
    # SLIDE 2 : MSP 지원 범위 (Support Scope)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "MSP 지원 범위 — 무엇을 요청받는가")
    tf = textbox(s, 0.6, 1.25, 12.1, 0.7); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "고객(예: KT)이 MSP에 요청하는 CycleCloud 지원 항목과, 각 항목이 매핑되는 실습 모듈입니다."
    p.font.size = Pt(14); p.font.color.rgb = DARK_TEXT
    card(s, 0.6, 2.1, 5.8, 4.8, "🎫 주요 지원 요청 항목", [
        "① CycleCloud 신규 생성 (서버·클러스터 구축)",
        "② 노드 증/감설 및 노드 사이즈(VM SKU) 변경",
        "③ Storage Account · Disk 마운트 구성",
        "④ CycleCloud 포털 사용 방법 안내",
        "⑤ 기본 트러블슈팅 (로그 확인 방법)",
    ])
    card(s, 6.8, 2.1, 5.9, 4.8, "📗 대응 실습 모듈 매핑", [
        "① → Module 03 (신규 생성/최초 구축)",
        "② → Module 04 (노드 증감설·사이즈)",
        "③ → Module 05 (스토리지·디스크 마운트)",
        "④ → Module 03 (포털 사용법)",
        "⑤ → Module 11 (트러블슈팅·로그 진단)",
        "심화: 06 cluster-init · 07 Accounting ·",
        "  08 사용자 · 09 파티션 · 10 GPU 모니터링",
    ], bg_color=GREEN_BG, border_color=GREEN_BORDER, title_color=GREEN_TITLE)

    # ---------------------------------------------------------------
    # SLIDE 3 : 책임 분담 모델 (Shared Responsibility)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "책임 분담 모델 — 누가 무엇을 담당하는가")
    card(s, 0.6, 1.4, 3.8, 5.5, "☁️ Azure (플랫폼)", [
        "IaaS 인프라 (VM/VMSS, 네트워크)",
        "Managed 서비스 (MySQL, Lustre,",
        "  Prometheus/Grafana, Storage)",
        "CycleCloud 제품·패치 (PG)",
        "GPU VM 용량(capa) 공급",
    ])
    card(s, 4.7, 1.4, 3.8, 5.5, "🛠️ Cloud MSP (운영·지원)", [
        "클러스터 생성·변경·삭제 대행",
        "노드 증/감설·사이즈 변경",
        "스토리지·디스크 마운트 구성",
        "포털/CLI 운영 및 사용 안내",
        "로그 확인·트러블슈팅·에스컬레이션",
    ], bg_color=CARD_BG, border_color=BLUE, title_color=BLUE)
    card(s, 8.8, 1.4, 3.9, 5.5, "🏢 고객 (워크로드 소유)", [
        "HPC/AI 작업(Job) 제출·실행",
        "애플리케이션·데이터 관리",
        "cluster-init 요구사항 정의",
        "용량·일정·예산 의사결정",
        "RI 등 상용 계약 소유",
    ])

    # ---------------------------------------------------------------
    # SLIDE 4 : 운영 대상 구성요소 (What MSP operates)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "MSP가 운영하는 CycleCloud 구성요소")
    tf = textbox(s, 0.6, 1.25, 12.1, 0.7); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CycleCloud는 스케줄러 클러스터를 오케스트레이션하는 관리 서비스입니다. MSP는 아래 3계층을 다룹니다."
    p.font.size = Pt(14); p.font.color.rgb = DARK_TEXT
    card(s, 0.6, 2.1, 3.8, 4.8, "① 컨트롤 플레인 (Server)", [
        "CycleCloud Server VM (cc-server)",
        "웹 포털 GUI + cycle_server 서비스",
        "cyclecloud CLI 설치 위치",
        "Managed Identity 로 노드 프로비저닝",
        "→ 로그·상태 확인의 시작점",
    ])
    card(s, 4.7, 2.1, 3.8, 4.8, "② 스케줄러 & 계산 노드", [
        "스케줄러 노드 (Slurm Master)",
        "  - slurmctld, azslurm CLI",
        "실행(Execute) 노드 = VMSS",
        "  - Job 수요 시 자동 증설",
        "  - 유휴 시 자동 감설",
    ])
    card(s, 8.8, 2.1, 3.9, 4.8, "③ 스토리지 & Locker", [
        "Blob Storage (Locker)",
        "  - 템플릿·cluster-init 스크립트",
        "공유 마운트 (NFS / Blob NFS)",
        "Azure Managed Lustre (고성능)",
        "노드 OS Disk",
    ])

    # ---------------------------------------------------------------
    # SLIDE 5 : 운영 인터페이스 & 접근 (How MSP operates)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "운영 인터페이스 & 서버 접근 방법")
    card(s, 0.6, 1.4, 3.8, 5.5, "🖥️ 웹 포털 (GUI)", [
        "https://<server-ip> 접속",
        "클러스터 생성/Edit/Start/Terminate",
        "노드 상태·Autoscale 설정",
        "Users·Accounting·Monitoring 탭",
        "→ 일상 운영의 기본 도구",
    ])
    card(s, 4.7, 1.4, 3.8, 5.5, "⌨️ CLI (cyclecloud · azslurm)", [
        "cc-server / 스케줄러 노드에 설치",
        "cyclecloud show_cluster / show_nodes",
        "cyclecloud start / terminate_cluster",
        "azslurm resume / suspend / scale",
        "→ 자동화·정밀 제어",
    ])
    card(s, 8.8, 1.4, 3.9, 5.5, "🔑 SSH 키 없는 접근", [
        "Azure Portal → cc-server →",
        "  '명령 실행(Run Command)'",
        "az vm run-command invoke ...",
        "Azure RBAC 권한만으로 로그 확인",
        "→ 키 분실·부재 시에도 진단 가능",
    ], bg_color=GREEN_BG, border_color=GREEN_BORDER, title_color=GREEN_TITLE)

    # ---------------------------------------------------------------
    # SLIDE 6 : 핵심 운영 작업 ① 생성 · 노드 증감설/사이즈
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "핵심 운영 작업 ① 클러스터 생성 & 노드 증/감설·사이즈")
    card(s, 0.6, 1.4, 5.8, 5.5, "🚀 신규 생성 (Module 03)", [
        "1. Locker용 Managed Identity·Storage 준비",
        "2. Marketplace 이미지로 CycleCloud VM 배포",
        "3. 설정 마법사: Site·License·관리자·구독 등록",
        "   (VM 시스템 ID에 Contributor +",
        "    Storage Blob Data Contributor 부여)",
        "4. Slurm 템플릿으로 클러스터 생성 & Start",
        "   Off→Acquiring→Preparing→Ready",
    ])
    card(s, 6.8, 1.4, 5.9, 5.5, "📐 노드 증/감설·사이즈 (Module 04)", [
        "증/감설: Autoscale Max Cores 조정 또는",
        "  azslurm resume/suspend 로 수동 제어",
        "사이즈 변경: 템플릿 VM SKU 변경 →",
        "  import_cluster --force → azslurm scale",
        "azslurm scale : 재시작 없이 변경 반영",
        "⚠️ 수동 azure.conf 편집은 scale 시 되돌아감",
        "⚠️ 노드 재시작/재생성은 최소화 (다음 장)",
    ])

    # ---------------------------------------------------------------
    # SLIDE 7 : 핵심 운영 작업 ② 스토리지 · 사용자
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "핵심 운영 작업 ② 스토리지 마운트 & 사용자 관리")
    card(s, 0.6, 1.4, 5.8, 5.5, "💾 스토리지·디스크 마운트 (Module 05)", [
        "공유 마운트 유형:",
        "  - Blob Storage NFS (데이터)",
        "  - Azure Files (홈/공유)",
        "  - Azure Managed Lustre (고성능 병렬)",
        "노드 OS Disk 크기 조정",
        "  ⚠️ BootDiskSize 변경 시 노드",
        "     Terminate/Start 필요",
        "마운트는 cluster-init 스크립트로 표준화",
    ])
    card(s, 6.8, 1.4, 5.9, 5.5, "👤 사용자 관리 (Module 08)", [
        "지원 방식: Built-in / AD / LDAP / Entra ID",
        "  (KT는 Built-in 사용 중)",
        "포털 Users 탭에서 계정 생성",
        "Keypair(pem) 등록으로 노드 SSH 접근",
        "Users 추가 시 jetpack 이 노드 OS 계정 생성",
        "sudo 필요 시 sudoers 반영",
    ])

    # ---------------------------------------------------------------
    # SLIDE 8 : 반드시 알아야 할 운영 주의사항 (Critical Gotchas)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "반드시 알아야 할 운영 주의사항 (KT 환경)")
    card(s, 0.6, 1.4, 12.133, 2.1, "⚠️ [최우선] 노드 재시작 최소화 원칙", [
        "GPU VM 용량(capa) 부족으로 재시작/재생성 시 재할당(Acquiring)이 실패할 수 있음",
        "KT GPU VM은 Reserved Instance(RI) 적용 → 의도치 않은 Scale-in 방지 필요",
        "조치: 재시작 불가피 시 반드시 Azure Capa 팀과 사전 일정 조율 후 진행",
    ], bg_color=WARN_BG, border_color=WARN_BORDER, title_color=ACCENT_ORANGE)
    card(s, 0.6, 3.7, 5.8, 3.2, "🛡️ Scale-in 방지 설정", [
        "Slurm 기본 Suspend time 300초(유휴 자동 감설)",
        "UI 'keep_alive' 옵션 오작동 사례 존재",
        "해결: slurm.conf 에 SuspendExcParts=hpc",
        "cluster-init 스크립트로 적용되어 있음",
    ])
    card(s, 6.8, 3.7, 5.9, 3.2, "🔁 cloud-init 금지 · cluster-init 사용", [
        "cloud-init 수정 시 노드 할당 실패:",
        "  'does not match ... attribute: CloudInit'",
        "  → 전체 클러스터 재기동 필요 (위험)",
        "cluster-init(converge)로 재부팅 없이 반영",
        "Locker에 4~5개 스크립트(scale-in·드라이버·마운트)",
    ], bg_color=WARN_BG, border_color=WARN_BORDER, title_color=ACCENT_ORANGE)

    # ---------------------------------------------------------------
    # SLIDE 9 : 트러블슈팅 & 로그 진단 (Module 11)
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "트러블슈팅 & 로그 진단 체계")
    card(s, 0.6, 1.4, 5.8, 5.5, "🔎 3단계 Triage 순서", [
        "① 포털 문제 노드의 Status 메시지 확인",
        "② 서버 로그 검색",
        "   /opt/cycle_server/logs/cycle_server.log",
        "③ 노드 내부 로그 확인",
        "   jetpack.log / cloud-init-output.log",
        "인프라 진단: az vm / vmss 상태, NSG, 쿼터",
    ])
    card(s, 6.8, 1.4, 5.9, 5.5, "🧰 자주 쓰는 진단 도구", [
        "sinfo -R : DRAIN/DOWN 노드 사유 확인",
        "azslurm retry_failed_nodes : 실패 노드 재시도",
        "systemctl restart slurmd : 좀비 노드 복구",
        "/opt/cycle/capture_logs.sh : 지원용 로그 번들",
        "healthagent : 불량 노드 자동 DRAIN(60s)",
        "→ 상세 절차는 Module 11 참조",
    ])

    # ---------------------------------------------------------------
    # SLIDE 10 : 에스컬레이션 & 버전 정보
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "Microsoft 에스컬레이션 & 필수 버전 정보")
    card(s, 0.6, 1.4, 5.8, 5.5, "📌 지원 시 필수 3종 버전", [
        "① CycleCloud 버전 (UI 우측 상단 '?')",
        "② Slurm 버전 (Cluster Edit > Advanced)",
        "③ jetpack(slurm) 프로젝트 버전 (노드 CLI)",
        "→ 버그/호환성/미지원 여부 판별의 첫 좌표",
        "→ 재현 환경 구성·업그레이드 경로 산정에 필수",
    ])
    card(s, 6.8, 1.4, 5.9, 5.5, "🗂️ KT 클러스터 버전 현황(참고)", [
        "Mexico Central:",
        "  CycleCloud 8.9.0-3754 / Slurm 25.11.4",
        "  jetpack slurm 4.0.8 · 학습128 + GPU4",
        "Korea Central:",
        "  CycleCloud 8.7.3-3438 / Slurm 23.11.19-1",
        "  jetpack slurm 3.0.12 · 학습6 + GPU1",
        "※ 시점에 따라 변동 — 지원 전 재확인",
    ], bg_color=GREEN_BG, border_color=GREEN_BORDER, title_color=GREEN_TITLE)

    # ---------------------------------------------------------------
    # SLIDE 11 : 온보딩 로드맵 & 실습 매뉴얼
    # ---------------------------------------------------------------
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "MSP 온보딩 로드맵 & 실습 매뉴얼")
    card(s, 0.6, 1.4, 3.8, 5.5, "📚 기본 (01~05)", [
        "01. 환경 개요·아키텍처",
        "02. 포털 사용 방법",
        "03. 신규 생성·최초 구축",
        "04. 노드 증/감설·사이즈",
        "05. 스토리지·디스크 마운트",
    ])
    card(s, 4.7, 1.4, 3.8, 5.5, "🚀 심화 (06~10)", [
        "06. cluster-init·커스텀 스크립트",
        "07. Slurm Job Accounting",
        "08. 사용자·키페어 관리",
        "09. 파티션 관리·추가",
        "10. 모니터링",
    ])
    card(s, 8.8, 1.4, 3.9, 5.5, "🔧 운용·런북 (11~12)", [
        "11. 트러블슈팅·로그 진단",
        "12. 데모 런북(진행 체크리스트)",
        "",
        "권장 순서: 매뉴얼 정독 →",
        "  실습 환경에서 01~12 핸즈온 →",
        "  지원 시나리오별 대응 숙달",
    ], bg_color=GREEN_BG, border_color=GREEN_BORDER, title_color=GREEN_TITLE)
    tf = textbox(s, 0.6, 6.95, 12.1, 0.5); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "상세 절차: 'Azure CycleCloud MSP 운영 매뉴얼(docs/ 12개 모듈)' 참조."
    p.font.size = Pt(12); p.font.italic = True; p.font.color.rgb = GRAY

    prs.save(output_path)
    print("Presentation saved.")


if __name__ == "__main__":
    build("docs/CycleCloud_MSP_소개.pptx")
